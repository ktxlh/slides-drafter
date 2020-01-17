import pptx
import config as c
import tools as t
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_PARAGRAPH_ALIGNMENT
import os
import argparse
import requests
import re

rules_path = os.path.join(c.root_path, "rules.txt")
ppt_bg_path = os.path.join(c.res_pictures, '2.jpg')
# laoluo_bg_path = os.path.join(c.res_pictures, '6.jpg')
# last_bg_path = os.path.join(c.res_pictures, '8.jpg')
# story_bg_path = os.path.join(c.res_pictures, '9.jpg')
ppt_file_name = os.path.join(c.outputs_documents_path, 'result.pptx')

default_path = os.path.join(c.res_pictures, 'moren.jpg')


# 厘米转英寸
def cm_to_in(cm):
    return Inches(cm / 2.54)

def dowmloadPic(html, keyword):
    pic_url = re.findall('"objURL":"(.*?)",', html, re.S)
    i = 1
    print('keyword:' + keyword + ' ' +'image...')
    for each in pic_url:
        print('downloading' + str(i) + '  '+'image，site:' + str(each))
        try:
            pic = requests.get(each, timeout=10)
        except:
            print('【Error】')
            continue

        dir = c.root_path + '/res/pictures/' + keyword  + '.jpg'
        fp = open(dir, 'wb')
        fp.write(pic.content)
        fp.close()
        if i==1:
            return


# 判断课件是否存在，不存在的新建一个空白
def ppt_existed(ppt_name):
    # if not os.path.exists(ppt_name):
    prs = Presentation()
    prs.slide_height = cm_to_in(14.35)
    prs.slide_width = cm_to_in(25.5)
    prs.save(ppt_name)


# 模板1：只有一个标题
def model_1(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(ppt_bg_path, cm_to_in(0), cm_to_in(0), cm_to_in(25.4), cm_to_in(14.288))
    title_box = slide.shapes.add_textbox(cm_to_in(3.89), cm_to_in(5.35), cm_to_in(17.61), cm_to_in(3.59))
    paragraph = title_box.text_frame.add_paragraph()
    paragraph.text = title
    paragraph.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    paragraph.font.size = Pt(60)
    paragraph.font.name = '微软雅黑'
    paragraph.font.color.rgb = RGBColor(255, 255, 255)


#  模板2_ori：有原文，有图片
def model_2_ori(prs, title, pic_path,txt):
    # title
    slide = prs.slides.add_slide(prs.slide_layouts[7])
    # slide.shapes.add_picture(ppt_bg_path, cm_to_in(0), cm_to_in(0), cm_to_in(25.4), cm_to_in(14.288))
    title_box_1 = slide.shapes.add_textbox(cm_to_in(2.27), cm_to_in(0.90), cm_to_in(2.24), cm_to_in(1.00))
    paragraph_1 = title_box_1.text_frame.add_paragraph()
    paragraph_1.text = title
    paragraph_1.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    paragraph_1.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    paragraph_1.font.size = Pt(44)
    paragraph_1.font.name = '微软雅黑'
    paragraph_1.font.color.rgb = RGBColor(0, 0, 0)
    #image
    try:
        img = slide.shapes.add_picture(pic_path, cm_to_in(0), cm_to_in(0), height=cm_to_in(5.72))
    except:
        img = slide.shapes.add_picture(default_path, cm_to_in(0), cm_to_in(0), height=cm_to_in(5.72))
    img.left = int(prs.slide_width / 2 + (prs.slide_width / 2 - img.width) / 2)
    img.top = int((prs.slide_height - img.height) / 2)
    body_shape = slide.shapes.placeholders
    body_shape[2].text = txt
#model_2_key1 :summary + one keyword
def model_2_key1(prs, title, pic_path,keyword):
    print("!")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    #title
    title_box_1 = slide.shapes.add_textbox(cm_to_in(2.27), cm_to_in(0.90), cm_to_in(2.24), cm_to_in(1.00))
    paragraph_1 = title_box_1.text_frame.add_paragraph()
    paragraph_1.text = title
    paragraph_1.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    paragraph_1.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    paragraph_1.font.size = Pt(32)
    paragraph_1.font.name = '微软雅黑'
    paragraph_1.font.color.rgb = RGBColor(0, 0, 0)
    #image
    try:
        img = slide.shapes.add_picture(pic_path, cm_to_in(0), cm_to_in(0), width=cm_to_in(10),height=cm_to_in(8.72))
    except:
        img = slide.shapes.add_picture(default_path, cm_to_in(0), cm_to_in(0), width=cm_to_in(10),height=cm_to_in(8.72))
    img.left = int(prs.slide_width / 3.5 - (prs.slide_width / 2 - img.width) / 2)
    img.top = int((prs.slide_height - img.height) / 2)
    title_box = slide.shapes.add_textbox(prs.slide_width / 3, cm_to_in(11.35), int(prs.slide_width / 3), cm_to_in(3.59))
    paragraph = title_box.text_frame.add_paragraph()
    paragraph.text = keyword
    paragraph.font.size = Pt(32)
    paragraph.font.name = '微软雅黑'
    paragraph.font.color.rgb = RGBColor(0, 0, 0)






#一大标题三张图三小标题
def model_4(prs, title, pic_path1,pic_path2,pic_path3,*content ):
    #title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # slide.shapes.add_picture(ppt_bg_path, cm_to_in(0), cm_to_in(0), cm_to_in(25.4), cm_to_in(14.288))
    title_box_1 = slide.shapes.add_textbox(cm_to_in(2.27), cm_to_in(0.90), cm_to_in(2.24), cm_to_in(1.00))
    paragraph_1 = title_box_1.text_frame.add_paragraph()
    paragraph_1.text = title
    paragraph_1.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    paragraph_1.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    paragraph_1.font.size = Pt(44)
    paragraph_1.font.name = '微软雅黑'
    paragraph_1.font.color.rgb = RGBColor(0, 0, 0)
    #subtitle
    module_width = (prs.slide_width - cm_to_in(1.27) * 2) / len(content)
    for i in range(0, 3):
        title_box = slide.shapes.add_textbox(cm_to_in(1.27) + i * module_width, cm_to_in(10.4), module_width,
                                             cm_to_in(2.39))
        paragraph = title_box.text_frame.add_paragraph()
        paragraph.text = content[i]
        paragraph.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        paragraph.font.size = Pt(32)
        paragraph.font.name = '微软雅黑'
        paragraph.font.color.rgb = RGBColor(0, 0, 0)
        #image
        try:
            img1 = slide.shapes.add_picture(pic_path1, cm_to_in(0), cm_to_in(0), width=cm_to_in(6),height=cm_to_in(5.72))
        except:
            img1 = slide.shapes.add_picture(default_path, cm_to_in(0), cm_to_in(0), width=cm_to_in(6),height=cm_to_in(5.72))
        img1.left = int(cm_to_in(1.27))
        img1.top = int((prs.slide_height - img1.height) / 2)
        try:
            img2 = slide.shapes.add_picture(pic_path2, cm_to_in(0), cm_to_in(0), width=cm_to_in(6),height=cm_to_in(5.72))
        except:
            img2 = slide.shapes.add_picture(default_path, cm_to_in(0), cm_to_in(0), width=cm_to_in(6),height=cm_to_in(5.72))
        img2.left = int(cm_to_in(9.27))
        img2.top = int((prs.slide_height - img2.height) / 2)
        try:
            img3 = slide.shapes.add_picture(pic_path3, cm_to_in(0), cm_to_in(0), width=cm_to_in(6),height=cm_to_in(5.72))
        except:
            img3 = slide.shapes.add_picture(default_path, cm_to_in(0), cm_to_in(0), width=cm_to_in(6),height=cm_to_in(5.72))
        img3.left = int(cm_to_in(17.27))
        img3.top = int((prs.slide_height - img3.height) / 2)
def model_5(prs, title, pic_path1,*content ):
    #title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # slide.shapes.add_picture(ppt_bg_path, cm_to_in(0), cm_to_in(0), cm_to_in(25.4), cm_to_in(14.288))
    title_box_1 = slide.shapes.add_textbox(cm_to_in(2.27), cm_to_in(0.90), cm_to_in(2.24), cm_to_in(1.00))
    paragraph_1 = title_box_1.text_frame.add_paragraph()
    paragraph_1.text = title
    paragraph_1.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    paragraph_1.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    paragraph_1.font.size = Pt(44)
    paragraph_1.font.name = '微软雅黑'
    paragraph_1.font.color.rgb = RGBColor(0, 0, 0)
    #subtitle
    module_width = (prs.slide_width - cm_to_in(1.27) * 2) / len(content)
    for i in range(0, 4):
        if i== 0 :
            title_box = slide.shapes.add_textbox(cm_to_in(1.27) , cm_to_in(3.4), module_width,
                                                 cm_to_in(2.39))
        elif i== 1 :
            title_box = slide.shapes.add_textbox(cm_to_in(1.27) + prs.slide_width/1.5, cm_to_in(3.4), module_width,
                                                 cm_to_in(2.39))
        elif i == 2:
            title_box = slide.shapes.add_textbox(cm_to_in(1.27) , cm_to_in(8.4), module_width,
                                                 cm_to_in(2.39))
        elif i== 3 :
            title_box = slide.shapes.add_textbox(cm_to_in(1.27) + prs.slide_width / 1.5,cm_to_in(8.4), module_width,cm_to_in(2.39))
        paragraph = title_box.text_frame.add_paragraph()
        paragraph.text = content[i]
        paragraph.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        paragraph.font.size = Pt(32)
        paragraph.font.name = '微软雅黑'
        paragraph.font.color.rgb = RGBColor(0, 0, 0)

        #image
        try:
            img1 = slide.shapes.add_picture(pic_path1, cm_to_in(0), cm_to_in(0), width=cm_to_in(6),height=cm_to_in(5.72))
        except:
            img1 = slide.shapes.add_picture(default_path, cm_to_in(0), cm_to_in(0), width=cm_to_in(6),height=cm_to_in(5.72))
        img1.left = int((prs.slide_width-img1.width)/2)
        img1.top = int((prs.slide_height - img1.height) / 2)


def model_8(prs, title, pic_path1,pic_path2 ,*content):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box_1 = slide.shapes.add_textbox(cm_to_in(0.27), cm_to_in(0.90), cm_to_in(5.24), cm_to_in(3.00))
    paragraph_1 = title_box_1.text_frame.add_paragraph()
    paragraph_1.text = title
    paragraph_1.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    paragraph_1.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    paragraph_1.font.size = Pt(44)
    paragraph_1.font.name = '微软雅黑'
    paragraph_1.font.color.rgb = RGBColor(0, 0, 0)
    # subtitle
    module_width = (prs.slide_width - cm_to_in(1.27) * 2) / len(content)
    for i in range(0, 2):
        title_box = slide.shapes.add_textbox(cm_to_in(1.27) + i * module_width, cm_to_in(10.4), module_width,
                                             cm_to_in(2.39))
        paragraph = title_box.text_frame.add_paragraph()
        paragraph.text = content[i]
        paragraph.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        paragraph.font.size = Pt(32)
        paragraph.font.name = '微软雅黑'
        paragraph.font.color.rgb = RGBColor(0, 0, 0)
        # image
        try:
            img1 = slide.shapes.add_picture(pic_path1, cm_to_in(0), cm_to_in(0), width=cm_to_in(7), height=cm_to_in(6.72))
        except:
            img1 = slide.shapes.add_picture(default_path, cm_to_in(0), cm_to_in(0), width=cm_to_in(7), height=cm_to_in(6.72))
        img1.left = int(cm_to_in(2.27))
        img1.top = int((prs.slide_height - img1.height) / 2)
        try:
            img2 = slide.shapes.add_picture(pic_path2, cm_to_in(0), cm_to_in(0), width=cm_to_in(7), height=cm_to_in(7.72))
        except:
            img2 = slide.shapes.add_picture(default_path, cm_to_in(0), cm_to_in(0), width=cm_to_in(7), height=cm_to_in(7.72))
        img2.left = int(cm_to_in(14.27))
        img2.top = int((prs.slide_height - img2.height) / 2)



# 读取配置文件调用模板的方法
def read_rules(prs, filename):
    if os.path.exists(filename):
        with open(filename, 'r+', encoding='utf-8') as f:
            for rule in f:
                word_list = rule.replace('\n', '').split('@')
                print('here')
                print(len(word_list))
                print(len(rule))
                if len(word_list) == 4:
                    word1 = word_list[1].strip()
                    word2 = word_list[2].strip()
                    word3 = word_list[3].strip()
                    url1 = 'http://image.baidu.com/search/flip?tn=baiduimage&ie=utf-8&word=' + word1 + '&ct=201326592&v=flip'
                    result1 = requests.get(url1)
                    dowmloadPic(result1.text, word1)
                    url2 = 'http://image.baidu.com/search/flip?tn=baiduimage&ie=utf-8&word=' + word2 + '&ct=201326592&v=flip'
                    result2 = requests.get(url2)
                    dowmloadPic(result2.text, word2)
                    url3 = 'http://image.baidu.com/search/flip?tn=baiduimage&ie=utf-8&word=' + word3 + '&ct=201326592&v=flip'
                    result3 = requests.get(url3)
                    dowmloadPic(result3.text, word3)
                    word1 = word_list[1].strip()
                    word2 = word_list[2].strip()
                    word3 = word_list[3].strip()
                    model_4(prs,word_list[0], os.path.join(c.res_pictures, word_list[1]+'.jpg'),os.path.join(c.res_pictures, word_list[2]+'.jpg'),os.path.join(c.res_pictures, word_list[3]+'.jpg'),word_list[1], word_list[2], word_list[3])
                if len(word_list) == 1:
                    model_1(prs, word_list[0])

                elif len(word_list)==2:
                    if len(rule)>20:
                        word=word_list[0]
                        url = 'http://image.baidu.com/search/flip?tn=baiduimage&ie=utf-8&word=' + word + '&ct=201326592&v=flip'
                        result = requests.get(url)
                        dowmloadPic(result.text, word)
                        # word_list 0:title    image is based on title not on keyword  1:original text(long)
                        model_2_ori(prs,word_list[0],os.path.join(c.res_pictures, word_list[0]+'.jpg'),word_list[1])
                    elif len(rule)>3:
                        word=word_list[1]
                        url = 'http://image.baidu.com/search/flip?tn=baiduimage&ie=utf-8&word=' + word + '&ct=201326592&v=flip'
                        result = requests.get(url)
                        dowmloadPic(result.text, word)
                        model_2_key1(prs, word_list[0], os.path.join(c.res_pictures, word_list[1] + '.jpg'),word_list[1])
                    else:
                        continue

                elif len(word_list)==3:
                    word1 =  word_list[1].strip()
                    word2 =  word_list[2].strip()
                    url1 = 'http://image.baidu.com/search/flip?tn=baiduimage&ie=utf-8&word=' + word1 + '&ct=201326592&v=flip'
                    result1 = requests.get(url1)
                    dowmloadPic(result1.text, word1)
                    url2 = 'http://image.baidu.com/search/flip?tn=baiduimage&ie=utf-8&word=' + word2 + '&ct=201326592&v=flip'
                    result2 = requests.get(url2)
                    dowmloadPic(result2.text, word2)
                    model_8(prs,word_list[0],os.path.join(c.res_pictures, word_list[1]+'.jpg'),os.path.join(c.res_pictures, word_list[2]+'.jpg'),word_list[1],word_list[2])
                elif len(word_list) == 5:
                    word0=word_list[0]
                    word1 = word_list[1]
                    word2 = word_list[2]
                    word3 = word_list[3]
                    word4 = word_list[4]
                    url0 = 'http://image.baidu.com/search/flip?tn=baiduimage&ie=utf-8&word=' + word0 + '&ct=201326592&v=flip'
                    result0 = requests.get(url0)
                    dowmloadPic(result0.text, word0)
                    url1 = 'http://image.baidu.com/search/flip?tn=baiduimage&ie=utf-8&word=' + word1 + '&ct=201326592&v=flip'
                    result1 = requests.get(url1)
                    dowmloadPic(result1.text, word1)
                    url2 = 'http://image.baidu.com/search/flip?tn=baiduimage&ie=utf-8&word=' + word2 + '&ct=201326592&v=flip'
                    result2 = requests.get(url2)
                    dowmloadPic(result2.text, word2)
                    url3 = 'http://image.baidu.com/search/flip?tn=baiduimage&ie=utf-8&word=' + word3 + '&ct=201326592&v=flip'
                    result3 = requests.get(url3)
                    dowmloadPic(result3.text, word3)
                    url4 = 'http://image.baidu.com/search/flip?tn=baiduimage&ie=utf-8&word=' + word4 + '&ct=201326592&v=flip'
                    result4 = requests.get(url4)
                    dowmloadPic(result4.text, word4)
                    model_5(prs, word_list[0], os.path.join(c.res_pictures, word_list[0] + '.jpg'),
                            word_list[1], word_list[2],
                            word_list[3],word_list[4])


if __name__ == '__main__':

    parser = argparse.ArgumentParser()
    parser.add_argument("-input_file", default='./rules.txt', help='rules_file')
    parser.add_argument('-output_file', default='./results.pptx', help='output_file')
    args = parser.parse_args()

    input_file = args.input_file
    output_file = args.output_file

    # t.is_dir_existed(c.outputs_documents_path)
    ppt_existed(output_file)
    presentation = Presentation(output_file)
    read_rules(presentation, input_file)
    presentation.save(output_file)

